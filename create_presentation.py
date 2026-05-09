#!/usr/bin/env python
"""
Generate professional PPT presentation for Model Performance & Metrics slide.
Creates a presentation-ready document with table comparison and speaker notes.
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call(['pip', 'install', 'python-pptx', '-q'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

import pandas as pd

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# ============================================================================
# SLIDE 1: Model Comparison Table
# ============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 244, 248)

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Model Performance Comparison"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.font.color.rgb = RGBColor(0, 51, 102)

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(9), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Binary Classification on Drop-Off Detection"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(18)
subtitle_para.font.color.rgb = RGBColor(100, 100, 100)

# Add table
rows, cols = 3, 7
left = Inches(0.5)
top = Inches(1.6)
width = Inches(9)
height = Inches(3.5)

table_shape = slide1.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Set column widths
col_widths = [Inches(2.0), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header row
headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "ROC-AUC", "Rank"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    # Style header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.bold = True
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(255, 255, 255)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)

# Data rows
data = [
    ["Logistic Regression", "91.4%", "88.8%", "90.9%", "89.9%", "97.4%", "🥇 1st"],
    ["Random Forest", "90.5%", "88.5%", "89.0%", "88.7%", "96.8%", "🥈 2nd"],
]

for row_idx, row_data in enumerate(data, 1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        # Style data cells
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                if row_idx == 1:  # Best model (Logistic Regression)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 102, 0)
        # Alternate row colors
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(230, 240, 250)

# Footer note
footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1.8))
footer_frame = footer_box.text_frame
footer_frame.word_wrap = True

footer_text = (
    "Key Insight: Both models exceed 90% accuracy threshold. "
    "Logistic Regression achieves the highest ROC-AUC (97.4%) and F1-Score (89.9%), "
    "making it the best model for this task. "
    "Class imbalance requires prioritizing Recall & F1-Score over raw Accuracy."
)
footer_frame.text = footer_text
footer_para = footer_frame.paragraphs[0]
footer_para.font.size = Pt(11)
footer_para.font.italic = True
footer_para.font.color.rgb = RGBColor(50, 50, 50)
footer_para.alignment = PP_ALIGN.LEFT

# ============================================================================
# SLIDE 2: Key Metrics & Talking Points
# ============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 244, 248)

# Title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame2 = title_box2.text_frame
title_frame2.text = "Best Model Performance & Strategic Talking Points"
title_para2 = title_frame2.paragraphs[0]
title_para2.font.size = Pt(36)
title_para2.font.bold = True
title_para2.font.color.rgb = RGBColor(0, 51, 102)

# Best Model Box
best_box = slide2.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(4.5), Inches(2.2))
best_box.fill.solid()
best_box.fill.fore_color.rgb = RGBColor(0, 102, 0)
best_box.line.color.rgb = RGBColor(0, 51, 0)
best_box.line.width = Pt(2)

best_text = best_box.text_frame
best_text.word_wrap = True
best_text.margin_top = Inches(0.15)
best_text.margin_left = Inches(0.15)

best_text.text = "🏆 BEST MODEL: LOGISTIC REGRESSION\n\nAccuracy: 91.4%\nPrecision: 88.8%\nRecall: 90.9%\nF1-Score: 89.9%\nROC-AUC: 97.4%"
for paragraph in best_text.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
    if paragraph.text.startswith("🏆"):
        for run in paragraph.runs:
            run.font.size = Pt(16)

# Talking Points
points_box = slide2.shapes.add_textbox(Inches(5.3), Inches(1.2), Inches(4.2), Inches(5.8))
points_frame = points_box.text_frame
points_frame.word_wrap = True

talking_points = [
    ("WHY F1 & RECALL MATTER", "Class imbalance in drop-off data means accuracy alone is misleading. Accuracy can be high even if the model ignores rare drop-offs (majority class bias)."),
    ("RECALL INTERPRETATION", "90.9% Recall means we catch ~91 out of 100 actual drop-offs. Missing the rest is costly (lost revenue, poor UX)."),
    ("F1-SCORE ADVANTAGE", "89.9% F1-Score balances precision & recall. It's the best metric when you care about both false positives AND false negatives."),
    ("BUSINESS IMPLICATION", "Better to flag a potential drop-off (false positive cost: $10) than miss a real one (false negative cost: $100)."),
]

for idx, (title, text) in enumerate(talking_points):
    if idx > 0:
        points_frame.add_paragraph()
    
    p = points_frame.paragraphs[-1 if idx > 0 else 0]
    p.text = f"• {title}"
    p.level = 0
    for run in p.runs:
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)
    
    points_frame.add_paragraph()
    p2 = points_frame.paragraphs[-1]
    p2.text = text
    p2.level = 1
    p2.space_before = Pt(3)
    p2.space_after = Pt(6)
    for run in p2.runs:
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(40, 40, 40)

# Save presentation
output_path = "results/Model_Performance_Metrics.pptx"
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"   - Slide 1: Model Comparison Table (both models >90% accuracy)")
print(f"   - Slide 2: Best Model Metrics + Strategic Talking Points")
